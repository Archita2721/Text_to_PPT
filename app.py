import os
import pptx
import base64
import traceback
import streamlit as st
import logging
from pptx.util import Pt
from openai import OpenAI
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Get OpenAI API key from environment variables
openai_api_key = os.getenv("OPENAI_API_KEY")

# Initialize OpenAI client
client = OpenAI(api_key=openai_api_key)

# Constants for font sizes
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

# Configure logging
logging.basicConfig(level=logging.INFO)


# Function to generate slide titles based on a topic using GPT-3.5-turbo
def generate_slide_titles(topic):
    prompt = f"Generate 5 slide titles for the topic '{topic}'."
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt},
        ],
    )
    logging.info(response.choices[0].message.content)
    return response.choices[0].message.content


# Function to generate slide content based on a slide title using GPT-3.5-turbo
def generate_slide_content(slide_title):
    prompt = f"Generate content for the slide: '{slide_title}'."
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt},
        ],
    )
    logging.info(response.choices[0].message.content)
    return response.choices[0].message.content


# Function to create a PowerPoint presentation given a topic, slide titles, and slide contents
def create_presentation(topic, slide_titles, slide_contents):
    try:
        # Initialize a PowerPoint presentation object
        prs = pptx.Presentation()
        # Choose the layout for slides
        slide_layout = prs.slide_layouts[1]

        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = topic

        # Add slides with titles and contents
        for slide_title, slide_content in zip(slide_titles, slide_contents):
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_title
            slide.shapes.placeholders[1].text = slide_content

            # Set font sizes for title and content
            slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = SLIDE_FONT_SIZE

        # Save the presentation
        ppt_filename = f"generated_ppt/{topic}_presentation.pptx"
        prs.save(ppt_filename)
        logging.info("Presentation saved successfully!")
    except Exception as e:
        # Log any errors that occur during presentation creation
        logging.error("Error occurred while creating presentation:")
        logging.error(traceback.format_exc())


# Main function to interact with the Streamlit interface
def main():
    # Streamlit app title
    st.title("PowerPoint Presentation Generator with GPT-3.5-turbo")

    # Text input for user to enter presentation topic
    topic = st.text_input("Enter the topic for your presentation:")

    # Button to trigger presentation generation
    generate_button = st.button("Generate Presentation")

    # Generate presentation when button is clicked and topic is provided
    if generate_button and topic:
        st.info("Generating presentation... Please wait.")
        # Generate slide titles
        slide_titles_string = generate_slide_titles(topic)
        slide_titles = slide_titles_string.split("\n")
        # Filter out empty slide titles
        filtered_slide_titles = [item for item in slide_titles if item.strip() != ""]
        logging.info("Slide Titles: %s", filtered_slide_titles)
        # Generate slide contents based on slide titles
        slide_contents = [
            generate_slide_content(title) for title in filtered_slide_titles
        ]
        logging.info("Slide Contents: %s", slide_contents)
        # Create the presentation
        create_presentation(topic, filtered_slide_titles, slide_contents)
        logging.info("Presentation generated successfully!")

        # Display success message and provide download link for the generated presentation
        st.success("Presentation generated successfully!")
        st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)


# Function to generate download link for the generated PowerPoint presentation
def get_ppt_download_link(topic):
    ppt_filename = f"generated_ppt/{topic}_presentation.pptx"
    ppt_path = os.path.join(os.getcwd(), ppt_filename)

    if not os.path.exists(ppt_path):
        # Log error if presentation file is not found
        logging.error(f"Error: Presentation file not found at {ppt_path}")
        return "Presentation file not found"

    with open(ppt_path, "rb") as file:
        ppt_contents = file.read()

    # Encode presentation file to base64 for download link
    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'


# Run the main function when the script is executed
if __name__ == "__main__":
    main()
